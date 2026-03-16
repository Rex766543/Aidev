#!/usr/bin/env python3
"""STRO Component Review - Consulting slide deck (Japanese)"""

import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# Colors
# ============================================================
PRIMARY_PURPLE   = RGBColor(0x70, 0x30, 0xA0)
HIGHLIGHT_PURPLE = RGBColor(0xCC, 0x00, 0xCC)
LIGHT_PURPLE     = RGBColor(0xEC, 0xE0, 0xF8)
DARK_TEXT        = RGBColor(0x26, 0x26, 0x26)
GRAY_TEXT        = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY       = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY         = RGBColor(0xCC, 0xCC, 0xCC)
COPYRIGHT_GRAY   = RGBColor(0x99, 0x99, 0x99)
WHITE            = RGBColor(0xFF, 0xFF, 0xFF)
AMBER            = RGBColor(0xFF, 0x80, 0x00)

SLIDE_W = 33.867
SLIDE_H = 19.05
COPYRIGHT_TEXT = "Confidential"
BREADCRUMBS = ["証明書アーキテクチャ", "ユースケース", "備考"]


# ============================================================
# Helpers
# ============================================================

def no_line(shape):
    shape.line.color.rgb = shape.fill.fore_color.rgb if shape.fill.type else WHITE


def add_common_elements(slide, slide_num, active_bc=None):
    # Slide number
    nb = slide.shapes.add_textbox(Cm(SLIDE_W - 1.5), Cm(SLIDE_H - 0.72), Cm(1.2), Cm(0.5))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(slide_num)
    r.font.size = Pt(10)
    r.font.name = "Meiryo"
    r.font.color.rgb = COPYRIGHT_GRAY

    # Copyright
    cb = slide.shapes.add_textbox(Cm(0.5), Cm(SLIDE_H - 0.72), Cm(SLIDE_W - 2.2), Cm(0.5))
    p = cb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = COPYRIGHT_TEXT
    r.font.size = Pt(8)
    r.font.name = "Meiryo"
    r.font.color.rgb = COPYRIGHT_GRAY

    # Breadcrumb tabs
    n = len(BREADCRUMBS)
    tab_h, tab_top, tab_gap = 0.52, 0.25, 0.07
    tab_w = 4.2
    start_x = SLIDE_W - n * tab_w - (n - 1) * tab_gap - 0.3

    for i, label in enumerate(BREADCRUMBS):
        x = start_x + i * (tab_w + tab_gap)
        is_active = (i == active_bc)
        shape = slide.shapes.add_shape(1, Cm(x), Cm(tab_top), Cm(tab_w), Cm(tab_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_PURPLE if is_active else LIGHT_GRAY
        shape.line.color.rgb = PRIMARY_PURPLE
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(9)
        r.font.name = "Meiryo"
        r.font.bold = is_active
        r.font.color.rgb = WHITE if is_active else DARK_TEXT


def add_title_area(slide, section_label, title, subtitle=""):
    # Section label (small, above title)
    if section_label:
        sb = slide.shapes.add_textbox(Cm(0.8), Cm(0.25), Cm(19.0), Cm(0.4))
        p = sb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = section_label
        r.font.size = Pt(11)
        r.font.name = "Meiryo"
        r.font.color.rgb = GRAY_TEXT

    # Main title
    tb = slide.shapes.add_textbox(Cm(0.8), Cm(0.7), Cm(19.5), Cm(1.1))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.name = "Meiryo"
    r.font.color.rgb = DARK_TEXT

    # Subtitle
    sub_y = 1.85
    if subtitle:
        sb2 = slide.shapes.add_textbox(Cm(0.8), Cm(sub_y), Cm(19.5), Cm(0.6))
        sb2.text_frame.word_wrap = True
        p = sb2.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.name = "Meiryo"
        r.font.color.rgb = GRAY_TEXT
        sub_y = 2.5

    # Separator line
    line = slide.shapes.add_shape(1, Cm(0.8), Cm(sub_y), Cm(SLIDE_W - 1.6), Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.color.rgb = LIGHT_GRAY

    return sub_y + 0.25  # content_top


def add_bullets(slide, items, left, top, width, height, base_size=11):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        text = item.get("text", "")
        level = item.get("level", 0)
        highlight = item.get("highlight", False)
        note = item.get("note", False)
        bold = item.get("bold", False)
        blank = item.get("blank", False)

        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if blank:
            r = para.add_run()
            r.text = ""
            r.font.size = Pt(5)
            continue

        indent = "    " * level
        bullet = "●" if level == 0 else "○"

        if bold and level == 0:
            color = PRIMARY_PURPLE
            size = base_size
            txt = text
            bullet_str = ""
        else:
            color = HIGHLIGHT_PURPLE if highlight else (GRAY_TEXT if note else DARK_TEXT)
            size = base_size - 1 if level > 0 else base_size
            if note:
                size = base_size - 2
            txt = text
            bullet_str = f"{indent}{bullet} "

        r = para.add_run()
        r.text = bullet_str + txt
        r.font.size = Pt(size)
        r.font.name = "Meiryo"
        r.font.color.rgb = color
        r.font.bold = bold


def col_header(slide, x, y, w, h, text, color=None):
    c = color or PRIMARY_PURPLE
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.color.rgb = c
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = "Meiryo"
    r.font.color.rgb = WHITE


def note_bar(slide, text, y=None):
    y = y or (SLIDE_H - 1.55)
    nb = slide.shapes.add_shape(1, Cm(0.8), Cm(y), Cm(SLIDE_W - 1.6), Cm(0.72))
    nb.fill.solid()
    nb.fill.fore_color.rgb = LIGHT_GRAY
    nb.line.color.rgb = MID_GRAY
    nb.line.width = Pt(0.5)
    txt = slide.shapes.add_textbox(Cm(1.1), Cm(y + 0.08), Cm(SLIDE_W - 2.0), Cm(0.55))
    txt.text_frame.word_wrap = True
    p = txt.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"※ {text}"
    r.font.size = Pt(9)
    r.font.name = "Meiryo"
    r.font.color.rgb = GRAY_TEXT


# ============================================================
# Slide 1: Title
# ============================================================
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Top purple band
    top = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(SLIDE_W), Cm(2.2))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY_PURPLE
    top.line.color.rgb = PRIMARY_PURPLE

    # Bottom purple band
    bot = slide.shapes.add_shape(1, Cm(0), Cm(SLIDE_H - 1.8), Cm(SLIDE_W), Cm(1.8))
    bot.fill.solid()
    bot.fill.fore_color.rgb = PRIMARY_PURPLE
    bot.line.color.rgb = PRIMARY_PURPLE

    # Left accent stripe
    acc = slide.shapes.add_shape(1, Cm(0), Cm(2.2), Cm(0.45), Cm(SLIDE_H - 4.0))
    acc.fill.solid()
    acc.fill.fore_color.rgb = HIGHLIGHT_PURPLE
    acc.line.color.rgb = HIGHLIGHT_PURPLE

    # Main title
    tb = slide.shapes.add_textbox(Cm(1.5), Cm(4.2), Cm(SLIDE_W - 3.0), Cm(3.5))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "STROコンポーネントレビュー"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.name = "Meiryo"
    r.font.color.rgb = DARK_TEXT

    # Subtitle
    sb = slide.shapes.add_textbox(Cm(1.5), Cm(8.3), Cm(SLIDE_W - 3.0), Cm(2.0))
    sb.text_frame.word_wrap = True
    p = sb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Nissan Charge / Plug & Charge サービス"
    r.font.size = Pt(20)
    r.font.name = "Meiryo"
    r.font.color.rgb = GRAY_TEXT

    # Badge: unofficial
    badge = slide.shapes.add_shape(1, Cm(1.5), Cm(11.0), Cm(4.5), Cm(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = AMBER
    badge.line.color.rgb = AMBER
    btf = badge.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "非公式版（Unofficial）"
    br.font.size = Pt(11)
    br.font.bold = True
    br.font.name = "Meiryo"
    br.font.color.rgb = WHITE

    # Footer date
    fb = slide.shapes.add_textbox(Cm(1.5), Cm(SLIDE_H - 1.4), Cm(SLIDE_W - 3.0), Cm(0.7))
    p = fb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Confidential"
    r.font.size = Pt(11)
    r.font.name = "Meiryo"
    r.font.color.rgb = WHITE


# ============================================================
# Slide 2: Background (MANDATORY)
# ============================================================
def slide_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 2, active_bc=None)
    ct = add_title_area(slide, "",
                        "背景・本日のレビュースコープ",
                        "今回のSTROコンポーネントレビューにおける前提と確認事項")

    cw = (SLIDE_W - 2.4) / 2
    ch = SLIDE_H - ct - 1.0

    cards = [
        {
            "title": "本レビューの位置づけ",
            "hdr_color": PRIMARY_PURPLE,
            "bg_color": LIGHT_PURPLE,
            "items": [
                {"text": "本STROはスプリント未リリースの非公式版"},
                {"text": "オープンポイント解消後、公式版リリース予定"},
                {"blank": True},
                {"text": "公式版にてEIMフォールバックの\nF/S結果・E/T削減に基づくアーキテクチャを策定"},
            ],
        },
        {
            "title": "本日のレビュースコープ",
            "hdr_color": GRAY_TEXT,
            "bg_color": LIGHT_GRAY,
            "items": [
                {"text": "NNAからの差分は少ない（EIMフォールバックを除く）"},
                {"text": "EIMフォールバック UCは今回スコープ外", "highlight": True},
                {"blank": True},
                {"text": "証明書アーキテクチャ（CPLC / V2G / multi-PKI）"},
                {"text": "UC01〜UC21 各UCの変更点確認"},
                {"text": "FBシート記入期限: 6月3日（火）", "bold": True},
            ],
        },
    ]

    for ci, card in enumerate(cards):
        x = 0.8 + ci * (cw + 0.4)
        col_header(slide, x, ct, cw, 0.72, card["title"], card["hdr_color"])

        bg = slide.shapes.add_shape(1, Cm(x), Cm(ct + 0.72), Cm(cw), Cm(ch - 0.72))
        bg.fill.solid()
        bg.fill.fore_color.rgb = card["bg_color"]
        bg.line.color.rgb = card["hdr_color"]
        bg.line.width = Pt(0.75)

        add_bullets(slide, card["items"], x + 0.25, ct + 0.95, cw - 0.5, ch - 1.1)


# ============================================================
# Slide 3: Certificate Architecture
# ============================================================
def slide_cert_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 3, active_bc=0)
    ct = add_title_area(slide, "証明書アーキテクチャ",
                        "CPLC・V2G Root証明書の概要",
                        "CPLCによるMissing Notification送信とmulti-PKI対応")

    cw = (SLIDE_W - 2.4) / 2
    ch = SLIDE_H - ct - 1.0

    # Left: CPLC
    col_header(slide, 0.8, ct, cw, 0.72, "CPLC Missing Notification")
    bg_l = slide.shapes.add_shape(1, Cm(0.8), Cm(ct + 0.72), Cm(cw), Cm(ch - 0.72))
    bg_l.fill.solid()
    bg_l.fill.fore_color.rgb = LIGHT_PURPLE
    bg_l.line.color.rgb = PRIMARY_PURPLE
    bg_l.line.width = Pt(0.75)
    add_bullets(slide, [
        {"text": "CPLCが未通知を検知しMissing Notificationを送信"},
        {"blank": True},
        {"text": "対象証明書"},
        {"text": "プロビジョニング証明書", "level": 1},
        {"text": "V2G Root証明書", "level": 1},
        {"blank": True},
        {"text": "送信経路: CPLC → IVC"},
    ], 1.05, ct + 0.95, cw - 0.5, ch - 1.2)

    # Right: multi-PKI
    rx = 0.8 + cw + 0.4
    col_header(slide, rx, ct, cw, 0.72, "V2G Root証明書 / multi-PKI")
    bg_r = slide.shapes.add_shape(1, Cm(rx), Cm(ct + 0.72), Cm(cw), Cm(ch - 0.72))
    bg_r.fill.solid()
    bg_r.fill.fore_color.rgb = LIGHT_PURPLE
    bg_r.line.color.rgb = PRIMARY_PURPLE
    bg_r.line.width = Pt(0.75)
    add_bullets(slide, [
        {"text": "EURリージョン: NNA同様のmulti-PKIアプローチ"},
        {"text": "V2G Root証明書は2枚構成"},
        {"blank": True},
        {"text": "1枚目: Hubject証明書（確定済み）"},
        {"text": "2枚目: TBD（ビジネス側で未確定）", "highlight": True},
        {"blank": True},
        {"text": "SOP時点での対応"},
        {"text": "Hubject証明書のみMissing Notificationを送信", "level": 1},
        {"text": "TBD証明書インストールはSOP時点では不要", "level": 1, "note": True},
    ], rx + 0.25, ct + 0.95, cw - 0.5, ch - 1.2)


# ============================================================
# Slide 4: V2G Cert Install Flow (2nd cert)
# ============================================================
def slide_v2g_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 4, active_bc=0)
    ct = add_title_area(slide, "証明書アーキテクチャ",
                        "V2G Root証明書（2枚目）インストールフロー",
                        "2枚目の証明書が確定した後の対応手順")

    steps = [
        {"num": "①", "label": "OTA\n設定更新",      "detail": "CPLC設定を\nOTA経由で更新"},
        {"num": "②", "label": "証明書\n2枚構成",     "detail": "V2G Root証明書\n2枚を設定"},
        {"num": "③", "label": "Missing\nNotif.",     "detail": "新Root証明書の\nNotificationを送信"},
        {"num": "④", "label": "Offboard\n検知",      "detail": "通知を検知"},
        {"num": "⑤", "label": "車両へ\nインストール", "detail": "証明書を\n車両に配信"},
    ]

    n = len(steps)
    flow_top = ct + 0.3
    flow_h = 2.3
    total_w = SLIDE_W - 1.6
    sw = total_w / n
    ov = 0.28

    for i, s in enumerate(steps):
        x = 0.8 + i * (sw - ov / 2)
        w = sw + (ov / 2 if i < n - 1 else -ov / 2)
        stype = 56 if i == 0 else 52
        shape = slide.shapes.add_shape(stype, Cm(x), Cm(flow_top), Cm(w), Cm(flow_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_PURPLE
        shape.line.color.rgb = PRIMARY_PURPLE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = s["num"]
        r1.font.size = Pt(9)
        r1.font.name = "Meiryo"
        r1.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = s["label"]
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.name = "Meiryo"
        r2.font.color.rgb = WHITE

        det = slide.shapes.add_textbox(
            Cm(x + 0.1), Cm(flow_top + flow_h + 0.2), Cm(w - 0.2), Cm(1.3))
        det.text_frame.word_wrap = True
        dp = det.text_frame.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        dr = dp.add_run()
        dr.text = s["detail"]
        dr.font.size = Pt(9)
        dr.font.name = "Meiryo"
        dr.font.color.rgb = GRAY_TEXT

    # Note box
    note_y = flow_top + flow_h + 2.0
    nb = slide.shapes.add_shape(1, Cm(0.8), Cm(note_y), Cm(SLIDE_W - 1.6), Cm(1.4))
    nb.fill.solid()
    nb.fill.fore_color.rgb = LIGHT_PURPLE
    nb.line.color.rgb = PRIMARY_PURPLE
    nb.line.width = Pt(0.75)
    ntxt = slide.shapes.add_textbox(Cm(1.1), Cm(note_y + 0.12), Cm(SLIDE_W - 2.2), Cm(1.1))
    ntf = ntxt.text_frame
    ntf.word_wrap = True
    np1 = ntf.paragraphs[0]
    nr1 = np1.add_run()
    nr1.text = "前提: 2枚目の証明書はビジネス側で確定待ち"
    nr1.font.size = Pt(10)
    nr1.font.bold = True
    nr1.font.name = "Meiryo"
    nr1.font.color.rgb = PRIMARY_PURPLE
    np2 = ntf.add_paragraph()
    nr2 = np2.add_run()
    nr2.text = "SOP時点ではHubject証明書のみ対応。詳細は公式STRO版に記載予定。"
    nr2.font.size = Pt(10)
    nr2.font.name = "Meiryo"
    nr2.font.color.rgb = DARK_TEXT


# ============================================================
# Slide 5: UC Overview
# ============================================================
def slide_uc_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 5, active_bc=1)
    ct = add_title_area(slide, "", "ユースケース一覧", "本レビュー対象ユースケース")

    uc_list = [
        ("UC01",   "プロビジョニング証明書インストール",      "NNA同等"),
        ("UC02",   "Nissan Charge 加入 & PnC 準備",          "コンポーネント変更あり"),
        ("UC03",   "充電器検索",                              "NNA同等"),
        ("UC04",   "PnC 充電",                               "一部変更あり"),
        ("UC05",   "充電履歴",                               "NNA同等"),
        ("UC06",   "PnC 無効化",                             "NNA同等"),
        ("UC10",   "契約解除",                               "NNA同等"),
        ("UC13",   "サービスプロモーション通知",              "NNA同等"),
        ("UC14",   "通知設定",                               "NNA同等"),
        ("UC15",   "支払情報更新",                           "I/F検討中"),
        ("UC17-20","RFID充電",                               "シーケンス仮定義"),
        ("UC21",   "アプリ充電",                             "I/F一部未確定"),
    ]

    cols, rows = 4, 3
    cw = (SLIDE_W - 1.6) / cols
    rh = (SLIDE_H - ct - 0.8) / rows
    changed_kw = ("変更", "検討", "未確定", "仮")

    for idx, (uc_id, title, status) in enumerate(uc_list):
        col = idx % cols
        row = idx // cols
        x = 0.8 + col * cw
        y = ct + row * rh

        card = slide.shapes.add_shape(1, Cm(x + 0.06), Cm(y + 0.06), Cm(cw - 0.12), Cm(rh - 0.12))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_PURPLE
        card.line.color.rgb = PRIMARY_PURPLE
        card.line.width = Pt(0.5)

        # UC badge
        bdg = slide.shapes.add_shape(1, Cm(x + 0.18), Cm(y + 0.18), Cm(2.3), Cm(0.55))
        bdg.fill.solid()
        bdg.fill.fore_color.rgb = PRIMARY_PURPLE
        bdg.line.color.rgb = PRIMARY_PURPLE
        btf = bdg.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = uc_id
        br.font.size = Pt(10)
        br.font.bold = True
        br.font.name = "Meiryo"
        br.font.color.rgb = WHITE

        # Title
        ttl = slide.shapes.add_textbox(Cm(x + 0.18), Cm(y + 0.85), Cm(cw - 0.36), Cm(0.7))
        ttl.text_frame.word_wrap = True
        tp = ttl.text_frame.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        tr.font.size = Pt(11)
        tr.font.bold = True
        tr.font.name = "Meiryo"
        tr.font.color.rgb = DARK_TEXT

        # Status
        is_changed = any(k in status for k in changed_kw)
        sc = slide.shapes.add_textbox(Cm(x + 0.18), Cm(y + rh - 0.58), Cm(cw - 0.36), Cm(0.45))
        sp2 = sc.text_frame.paragraphs[0]
        sr = sp2.add_run()
        sr.text = f"  {status}"
        sr.font.size = Pt(9)
        sr.font.name = "Meiryo"
        sr.font.color.rgb = HIGHLIGHT_PURPLE if is_changed else GRAY_TEXT


# ============================================================
# Slide factory: standard UC slide
# ============================================================
def std_slide(prs, slide_num, uc_label, title, subtitle, items,
              right_items=None, note=None, active_bc=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, slide_num, active_bc=active_bc)
    ct = add_title_area(slide, uc_label, title, subtitle)

    note_h = 0.85 if note else 0.0
    ch = SLIDE_H - ct - 0.75 - note_h

    if right_items is not None:
        cw = (SLIDE_W - 2.4) / 2
        lhdr = right_items[0].get("col_header", "") if isinstance(right_items[0], dict) else ""
        rhdr = items[0].get("col_header", "") if isinstance(items[0], dict) else ""

        # Two columns
        lx, rx = 0.8, 0.8 + cw + 0.4
        if lhdr or rhdr:
            col_header(slide, lx, ct, cw, 0.6, lhdr)
            col_header(slide, rx, ct, cw, 0.6, rhdr, GRAY_TEXT)
            ct2 = ct + 0.68
            ch2 = ch - 0.68
        else:
            ct2, ch2 = ct, ch

        add_bullets(slide, items, lx, ct2, cw, ch2)
        add_bullets(slide, right_items, rx, ct2, cw, ch2)
    else:
        add_bullets(slide, items, 0.8, ct, SLIDE_W - 1.6, ch)

    if note:
        note_bar(slide, note)
    return slide


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    # 1. Title
    slide_title(prs)

    # 2. Background (mandatory)
    slide_background(prs)

    # 3. Certificate overview
    slide_cert_overview(prs)

    # 4. V2G install flow
    slide_v2g_flow(prs)

    # 5. UC overview
    slide_uc_overview(prs)

    # 6. UC01
    std_slide(prs, 6, "UC01",
              "プロビジョニング証明書インストール",
              "KMRによる証明書オンボードインストール",
              [
                  {"text": "KMRが各証明書をオンボードシステムにインストール"},
                  {"text": "処理の詳細はNNAと同等"},
                  {"blank": True},
                  {"text": "留意事項", "bold": True},
                  {"text": "一部プロセスは最新LE1実装状況を未反映", "note": True, "level": 1},
                  {"text": "最新情報は公式STRO版にて反映予定", "note": True, "level": 1},
              ])

    # 7. UC02: Overview
    std_slide(prs, 7, "UC02",
              "Nissan Charge 加入 & PnC 準備",
              "コンポーネント変更点と全体構成",
              [
                  {"text": "コンポーネント変更点", "bold": True},
                  {"text": "通知: Batch.com → FCM（Firebase Cloud Messaging）",
                   "level": 1, "highlight": True},
                  {"text": "NRG（Nissan Regional Gateway）: Global IS方針により削除",
                   "level": 1, "highlight": True},
                  {"blank": True},
                  {"text": "全体構成（3パート）", "bold": True},
                  {"text": "① SID適用性確認", "level": 1},
                  {"text": "② ユーザー登録プロセス", "level": 1},
                  {"text": "③ PnC契約証明書生成", "level": 1},
              ])

    # 8. UC02: SID Check
    std_slide(prs, 8, "UC02",
              "SID適用性確認 & PnCフラグ",
              "Nissan Storeページ表示時のチェックフロー",
              [
                  {"text": "Nissan Storeサービスページ表示時にBSPがSID適用性を確認"},
                  {"text": "PnC適用性確認: EMPがKMR・HubjectでPCIDの有効化を検証"},
                  {"text": "→ UC01（プロビジョニング証明書インストール完了）の確認に相当"},
                  {"blank": True},
                  {"text": "PnC適用性フラグ", "bold": True},
                  {"text": "false: PnCなしでNissan Chargeを初回起動", "level": 1},
                  {"text": "true :  PnCありで2回目以降の起動", "level": 1},
              ])

    # 9. UC02: Enrollment & Contract Cert
    std_slide(prs, 9, "UC02",
              "ユーザー登録 & 契約証明書生成",
              "加入フローと証明書生成処理",
              [
                  {"text": "ユーザー登録フロー", "bold": True},
                  {"text": "「Subscribe Nissan Charge」タップ → キャッシュ削除（料金プラン・T&C）",
                   "level": 1},
                  {"text": "料金プラン選択 → T&C同意 → 支払情報登録 → UID有効化", "level": 1},
                  {"text": "途中離脱時にプロセスを再開始させる設計（キャッシュ削除→再取得）",
                   "level": 1, "note": True},
                  {"blank": True},
                  {"text": "契約証明書生成（PnCフラグ=true の場合）", "bold": True},
                  {"text": "EMPがE-Mobilityアカウントを作成", "level": 1},
                  {"text": "HubjectへContractCertificate生成をリクエスト", "level": 1},
                  {"text": "以降はNNAアーキテクチャと同等", "level": 1},
                  {"text": "インストール失敗時: アプリが応答を受け取るまでEMPがリトライ",
                   "level": 1, "note": True},
              ])

    # 10. UC03
    std_slide(prs, 10, "UC03",
              "充電器検索",
              "EPoIリスト同期とマップ表示",
              [
                  {"text": "EMPがEPoIリストをERPと同期"},
                  {"text": "アプリがEMPからEPoIリストを取得しマップに充電器を表示"},
                  {"blank": True},
                  {"text": "共通前処理", "bold": True},
                  {"text": "オーナーユーザー確認（PnCサービスはオーナーのみ利用可）", "level": 1},
                  {"text": "Nissan Charge加入状況確認", "level": 1},
                  {"blank": True},
                  {"text": "マップ表示フロー", "bold": True},
                  {"text": "Googleマップ取得 → BFF経由でEMPにリクエスト（場所・フィルタ条件）",
                   "level": 1},
                  {"text": "EMP → 充電ステーションID・充電器ID・リアルタイム空き情報を返却",
                   "level": 1},
                  {"text": "アプリがマップ上に充電器情報を重ねて表示", "level": 1},
              ])

    # 11. UC04
    std_slide(prs, 11, "UC04",
              "PnC 充電",
              "全体構成はNNAからほぼ変更なし",
              [
                  {"text": "全体アーキテクチャ: NNAと同等"},
                  {"blank": True},
                  {"text": "変更点", "bold": True},
                  {"text": "KMR通知ハンドリングの変更", "level": 1, "highlight": True},
                  {"text": "SVT（盗難車両追跡）ステータス確認箇所の変更", "level": 1, "highlight": True},
              ])

    # 12. UC05 + UC06 (two-column)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 12, active_bc=1)
    ct = add_title_area(slide, "UC05 / UC06", "充電履歴 / PnC 無効化",
                        "履歴参照と無効化フロー")
    cw = (SLIDE_W - 2.4) / 2
    ch = SLIDE_H - ct - 0.75

    col_header(slide, 0.8, ct, cw, 0.65, "UC05: 充電履歴")
    add_bullets(slide, [
        {"text": "ユーザー認証後"},
        {"text": "充電履歴リストを取得", "level": 1},
        {"text": "個別の充電記録詳細を取得", "level": 1},
    ], 0.8, ct + 0.73, cw, ch - 0.75)

    rx = 0.8 + cw + 0.4
    col_header(slide, rx, ct, cw, 0.65, "UC06: PnC 無効化")
    add_bullets(slide, [
        {"text": "加入状況確認 → PnC充電状況確認 → 無効化リクエスト送信"},
        {"text": "KMRからvNextへリクエスト時: Appへサイレントプッシュ通知", "level": 1},
        {"blank": True},
        {"text": "オンボードステータスがdisabledへ変更後"},
        {"text": "CPLCがコマンドレスポンス返却 → EMPがPnCサービス状態更新", "level": 1},
        {"text": "アプリが完了通知を受信", "level": 1},
        {"blank": True},
        {"text": "暫定通知 (Interim Notification)", "bold": True},
        {"text": "アクションステータス更新 → Appへサイレントプッシュ", "level": 1},
        {"text": "サイレントプッシュ失敗時のためポーリングループを実装",
         "level": 1, "note": True},
    ], rx, ct + 0.73, cw, ch - 0.75)

    # 13. UC10
    std_slide(prs, 13, "UC10",
              "契約解除",
              "通常解除フローと複数トリガー対応",
              [
                  {"text": "複数のトリガー（ユーザー申請・期限切れ等）に対応"},
                  {"blank": True},
                  {"text": "通常解除フロー", "bold": True},
                  {"text": "ユーザーがアプリから解除申請 → BSPへリクエスト送信", "level": 1},
                  {"text": "UID無効化", "level": 1},
                  {"text": "EMPへ契約削除リクエスト送信", "level": 1},
                  {"text": "HubjectがContractCertificateを削除", "level": 1},
                  {"text": "Hubject → EMP・KMRへ通知 → オンボード証明書削除", "level": 1},
                  {"text": "EMPの顧客データを削除", "level": 1},
                  {"text": "ユーザーへ解除完了通知を送信", "level": 1},
              ])

    # 14. UC13 + UC14 (two-column)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_common_elements(slide, 14, active_bc=1)
    ct = add_title_area(slide, "UC13 / UC14", "サービス通知 & 通知設定",
                        "プロモーション通知の配信と設定管理")
    cw = (SLIDE_W - 2.4) / 2
    ch = SLIDE_H - ct - 0.75

    col_header(slide, 0.8, ct, cw, 0.65, "UC13: サービスプロモーション通知")
    add_bullets(slide, [
        {"text": "対象通知"},
        {"text": "Nissan Chargeサービス開始通知", "level": 1},
        {"text": "PnCサービス開始通知", "level": 1},
        {"blank": True},
        {"text": "開始タイミングにより2種類の通知が存在"},
        {"text": "通知内容は支払情報登録状況によって異なる", "note": True},
    ], 0.8, ct + 0.73, cw, ch - 0.75)

    rx = 0.8 + cw + 0.4
    col_header(slide, rx, ct, cw, 0.65, "UC14: 通知設定")
    add_bullets(slide, [
        {"text": "ユーザーがプロモーション通知設定を変更可能"},
        {"text": "KMR通知設定とは独立した設定"},
        {"text": "通知画面から直接プロモーション通知をオフにできる"},
    ], rx, ct + 0.73, cw, ch - 0.75)

    # 15. UC15
    std_slide(prs, 15, "UC15",
              "支払情報更新",
              "Nissan Storeでの支払情報変更フロー",
              [
                  {"text": "ユーザーがNissan Storeで支払情報を更新"},
                  {"blank": True},
                  {"text": "フロー", "bold": True},
                  {"text": "BSPで支払情報を更新", "level": 1},
                  {"text": "EMP経由でERPへ情報を転送", "level": 1},
                  {"text": "請求処理に使用", "level": 1},
                  {"blank": True},
                  {"text": "保存方法はOctopus・ISインターフェースチームと検討中", "highlight": True},
              ])

    # 16. UC17-20
    std_slide(prs, 16, "UC17-20",
              "RFID 充電",
              "EMP・ERP・R&Dコンポーネント間インターフェース",
              [
                  {"text": "RFID充電プロセスに関するユースケース群"},
                  {"blank": True},
                  {"text": "現状", "bold": True},
                  {"text": "EMP・ERP・R&Dコンポーネント間のI/Fは現在検討中", "highlight": True, "level": 1},
                  {"text": "現行シーケンス図は現時点の理解に基づく仮定義", "note": True, "level": 1},
                  {"blank": True},
                  {"text": "確定次第シーケンス図を更新予定"},
              ],
              note="シーケンス図は現時点での仮定義。I/F確定後に更新。")

    # 17. UC21
    std_slide(prs, 17, "UC21",
              "アプリ充電",
              "アプリ操作による充電起動シーケンス",
              [
                  {"text": "アプリベースの充電操作シーケンス"},
                  {"text": "ラフアーキテクチャからの大きな変更なし"},
                  {"blank": True},
                  {"text": "留意事項", "bold": True},
                  {"text": "関連I/FのほとんどがISコンポーネントに関連", "level": 1},
                  {"text": "一部インターフェースは未確定", "highlight": True, "level": 1},
                  {"text": "現シーケンス図は暫定仮定義として扱うこと", "note": True, "level": 1},
              ],
              active_bc=2,
              note="ISコンポーネントのI/F確定後にシーケンス図を更新予定。")

    # Save
    out_dir = "/Users/ishikawaryuusuke/Desktop/aidev/presentations"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "STRO_component_review.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
