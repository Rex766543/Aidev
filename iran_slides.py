from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# カラー定義
COLOR_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 濃紺（背景）
COLOR_BG_CARD   = RGBColor(0x16, 0x21, 0x3E)   # やや明るい紺（カード）
COLOR_ACCENT    = RGBColor(0xE8, 0x3A, 0x3A)   # 赤（アクセント）
COLOR_GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # ゴールド
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xCF, 0xD8, 0xE6)   # 薄青白

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # 完全空白


def add_bg(slide, color=COLOR_BG_DARK):
    """スライド全体の背景矩形を追加"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0,
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=COLOR_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Meiryo"
    return txBox


def add_bullet_box(slide, left, top, width, height, items, font_size=14,
                   bullet="●", color=COLOR_WHITE, spacing=0.4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing * font_size)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Meiryo"
    return txBox


# ============================================================
# スライド1: タイトル
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1)

# 左側アクセントバー
add_rect(slide1,
         Inches(0), Inches(0),
         Inches(0.35), prs.slide_height,
         COLOR_ACCENT)

# ゴールドライン（上部装飾）
add_rect(slide1,
         Inches(0.35), Inches(2.6),
         Inches(12.98), Pt(3),
         COLOR_GOLD)

# メインタイトル
add_textbox(slide1,
            Inches(0.7), Inches(2.75),
            Inches(12), Inches(1.5),
            "イラン情勢まとめ", 52, bold=True,
            color=COLOR_WHITE, align=PP_ALIGN.LEFT)

# サブタイトル
add_textbox(slide1,
            Inches(0.7), Inches(3.65),
            Inches(10), Inches(0.8),
            "2026年3月現在 ／ 軍事衝突・国内混乱・国際社会の対応", 18,
            color=COLOR_LIGHT, align=PP_ALIGN.LEFT)

# 日付
add_textbox(slide1,
            Inches(0.7), Inches(6.5),
            Inches(4), Inches(0.5),
            "2026年3月3日", 13,
            color=COLOR_GOLD, align=PP_ALIGN.LEFT)

# ============================================================
# スライド2: 国内反政府デモ
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2)
add_rect(slide2, Inches(0), Inches(0), Inches(0.35), prs.slide_height, COLOR_ACCENT)

# ヘッダー帯
add_rect(slide2, Inches(0.35), Inches(0), Inches(12.98), Inches(1.45), COLOR_BG_CARD)

add_textbox(slide2,
            Inches(0.6), Inches(0.2),
            Inches(12), Inches(0.9),
            "01  ／  国内反政府デモの勃発（2025年12月〜2026年1月）",
            20, bold=True, color=COLOR_GOLD)

# 左カード
add_rect(slide2, Inches(0.6), Inches(1.65), Inches(5.6), Inches(5.3), COLOR_BG_CARD)
add_textbox(slide2,
            Inches(0.75), Inches(1.75),
            Inches(5.3), Inches(0.55),
            "背景と経緯", 15, bold=True, color=COLOR_ACCENT)
add_bullet_box(slide2,
               Inches(0.75), Inches(2.35),
               Inches(5.3), Inches(4.3),
               [
                   "インフレ・食料品価格高騰・通貨暴落への不満が爆発",
                   "2025年12月28日、複数都市で大規模デモ勃発",
                   "商店主・大学生・労働者・女性・少数民族が参加",
                   "「経済改善」から「政権打倒」へ要求が急速に拡大",
                   "地方都市を含む全国規模へと波及",
               ],
               font_size=13.5, bullet="▶")

# 右カード
add_rect(slide2, Inches(6.5), Inches(1.65), Inches(6.2), Inches(5.3), COLOR_BG_CARD)
add_textbox(slide2,
            Inches(6.65), Inches(1.75),
            Inches(5.8), Inches(0.55),
            "政府の対応", 15, bold=True, color=COLOR_ACCENT)
add_bullet_box(slide2,
               Inches(6.65), Inches(2.35),
               Inches(5.8), Inches(4.3),
               [
                   "治安部隊が実弾発砲、数十人が逮捕",
                   "イラクのシーア派民兵800人を鎮圧に動員",
                   "インターネット検閲を強化・情報統制を徹底",
                   "地方での実態把握が国際社会にとり困難に",
                   "トランプ米大統領がデモ参加者に声援",
               ],
               font_size=13.5, bullet="▶")

# ============================================================
# スライド3: 米・イスラエルによる軍事攻撃
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3)
add_rect(slide3, Inches(0), Inches(0), Inches(0.35), prs.slide_height, COLOR_ACCENT)

add_rect(slide3, Inches(0.35), Inches(0), Inches(12.98), Inches(1.45), COLOR_BG_CARD)
add_textbox(slide3,
            Inches(0.6), Inches(0.2),
            Inches(12), Inches(0.9),
            "02  ／  米・イスラエルによる大規模軍事攻撃（2026年2月28日）",
            20, bold=True, color=COLOR_GOLD)

# 3カラムカード
for i, (title, bullets) in enumerate([
    ("攻撃の概要", [
        "2026年2月28日に攻撃開始",
        "約500か所の標的を一斉攻撃",
        "防空システム・ミサイル発射装置",
        "核施設・弾道ミサイルインフラ",
        "指導部の拠点・軍事司令部",
    ]),
    ("作戦の目的", [
        "イランの政権交代を誘導",
        "核開発プログラムの無力化",
        "地域の安全保障上の脅威除去",
        "IAEAは核施設への直撃を未確認",
        "米兵3名死亡・5名重傷が発生",
    ]),
    ("使用された手段", [
        "大量の航空機による波状攻撃",
        "精密誘導爆弾・巡航ミサイル",
        "米・イスラエルの合同作戦",
        "事前の外交交渉は不成立",
        "湾岸基地から出撃",
    ]),
]):
    col_left = Inches(0.6 + i * 4.25)
    add_rect(slide3, col_left, Inches(1.65), Inches(3.95), Inches(5.3), COLOR_BG_CARD)
    add_textbox(slide3,
                col_left + Inches(0.15), Inches(1.75),
                Inches(3.7), Inches(0.55),
                title, 15, bold=True, color=COLOR_ACCENT)
    add_bullet_box(slide3,
                   col_left + Inches(0.15), Inches(2.35),
                   Inches(3.7), Inches(4.3),
                   bullets, font_size=13, bullet="▶")

# ============================================================
# スライド4: ハメネイ死亡とイランの報復
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4)
add_rect(slide4, Inches(0), Inches(0), Inches(0.35), prs.slide_height, COLOR_ACCENT)

add_rect(slide4, Inches(0.35), Inches(0), Inches(12.98), Inches(1.45), COLOR_BG_CARD)
add_textbox(slide4,
            Inches(0.6), Inches(0.2),
            Inches(12), Inches(0.9),
            "03  ／  最高指導者ハメネイ師の死亡とイランの報復",
            20, bold=True, color=COLOR_GOLD)

# 重要事実ボックス（上段）
add_rect(slide4, Inches(0.6), Inches(1.65), Inches(11.9), Inches(1.05), RGBColor(0x4A, 0x0A, 0x0A))
add_rect(slide4, Inches(0.6), Inches(1.65), Inches(0.08), Inches(1.05), COLOR_ACCENT)
add_textbox(slide4,
            Inches(0.85), Inches(1.75),
            Inches(11.4), Inches(0.8),
            "2026年3月1日：最高指導者アリ・ハメネイ師（享年86歳）が攻撃で死亡。娘・息子の配偶者・孫も同時に死亡。イラン国営メディアが確認。",
            14.5, bold=True, color=COLOR_WHITE)

# 左：報復の内容
add_rect(slide4, Inches(0.6), Inches(2.85), Inches(5.6), Inches(4.1), COLOR_BG_CARD)
add_textbox(slide4,
            Inches(0.75), Inches(2.95),
            Inches(5.3), Inches(0.55),
            "イランの報復攻撃", 15, bold=True, color=COLOR_ACCENT)
add_bullet_box(slide4,
               Inches(0.75), Inches(3.55),
               Inches(5.3), Inches(3.2),
               [
                   "イスラエル・米軍駐留国に向けミサイル発射",
                   "標的：サウジアラビア・UAE・カタール・バーレーン・クウェート・イラク・ヨルダン",
                   "大半が迎撃されたが一部が着弾",
                   "ヒズボラがイスラエル・ハイファのIDF基地攻撃",
                   "在リヤド米国大使館がドローン攻撃を受ける",
               ],
               font_size=13, bullet="▶")

# 右：イラン政府の姿勢
add_rect(slide4, Inches(6.5), Inches(2.85), Inches(6.2), Inches(4.1), COLOR_BG_CARD)
add_textbox(slide4,
            Inches(6.65), Inches(2.95),
            Inches(5.8), Inches(0.55),
            "イラン政府の姿勢", 15, bold=True, color=COLOR_ACCENT)
add_bullet_box(slide4,
               Inches(6.65), Inches(3.55),
               Inches(5.8), Inches(3.2),
               [
                   "「立場は揺るがない」と国営メディアが声明",
                   "「死はイランをより強くした」と強調",
                   "40日間の喪に服すると宣言・7日間の国民休暇",
                   "イスラエルは引き続きテヘランの数十か所を攻撃",
                   "最高指導者の後継問題が今後の焦点に",
               ],
               font_size=13, bullet="▶")

# ============================================================
# スライド5: 国際社会の反応と今後の影響
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5)
add_rect(slide5, Inches(0), Inches(0), Inches(0.35), prs.slide_height, COLOR_ACCENT)

add_rect(slide5, Inches(0.35), Inches(0), Inches(12.98), Inches(1.45), COLOR_BG_CARD)
add_textbox(slide5,
            Inches(0.6), Inches(0.2),
            Inches(12), Inches(0.9),
            "04  ／  国際社会の反応と今後の影響",
            20, bold=True, color=COLOR_GOLD)

# 上段：各国反応（3列）
for i, (country, stance, detail) in enumerate([
    ("中国・ロシア", "強く非難", "「国連憲章違反」（中国）\n「不当な武力行使」（ロシア）\n両国ともイラン支持を表明"),
    ("EU・国際社会", "自制を要求", "「最大限の自制を求める」\n民間人保護・国際法の遵守を要求\n外交的解決を呼びかけ"),
    ("米・トランプ政権", "継続の意向", "「米兵の死を報復する」と宣言\n作戦継続を表明\n以前よりイランへの軍事攻撃を示唆"),
]):
    col_left = Inches(0.6 + i * 4.25)
    add_rect(slide5, col_left, Inches(1.65), Inches(3.95), Inches(2.4), COLOR_BG_CARD)
    add_textbox(slide5, col_left + Inches(0.15), Inches(1.75),
                Inches(3.7), Inches(0.5), country, 14, bold=True, color=COLOR_GOLD)
    add_textbox(slide5, col_left + Inches(0.15), Inches(2.2),
                Inches(3.7), Inches(0.45), stance, 13, bold=True, color=COLOR_ACCENT)
    add_textbox(slide5, col_left + Inches(0.15), Inches(2.62),
                Inches(3.7), Inches(1.3), detail, 12, color=COLOR_LIGHT)

# 下段：経済・社会的影響
add_rect(slide5, Inches(0.6), Inches(4.2), Inches(11.9), Inches(2.75), COLOR_BG_CARD)
add_textbox(slide5,
            Inches(0.75), Inches(4.3),
            Inches(5), Inches(0.5),
            "経済・社会・日本への影響", 15, bold=True, color=COLOR_ACCENT)

impacts = [
    "✈  湾岸航空各社（エミレーツ・エティハド・カタール航空など）が運航停止→中東の航空ハブ機能がまひ",
    "🛢  ホルムズ海峡・紅海の緊張が頂点に達し、原油価格が急騰。世界のサプライチェーンに深刻な影響",
    "🇯🇵  日本：高市首相が補正予算編成の可能性に言及。エネルギー安全保障の見直しが急務",
    "🌍  世界各地で反戦デモが拡大。イラン支持・攻撃非難の声が高まる",
]
for j, text in enumerate(impacts):
    add_textbox(slide5,
                Inches(0.75), Inches(4.8 + j * 0.47),
                Inches(11.5), Inches(0.45),
                text, 12.5, color=COLOR_LIGHT)

output_path = "/Users/ishikawaryuusuke/Desktop/aidev/iran_situation_2026.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")
